"""Generate presentation PPT for CS 513 Final Project."""
import os
import json
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

os.chdir(os.path.dirname(os.path.abspath(__file__)))

RESULTS_DIR = "../Results"
FIGURES_DIR = "../Results/Figures"
A_FIGURES = "../../Version_A/Results/Figures"
A_RESULTS = "../../Version_A/Results"
B_RESULTS = "../../Version_B"
OUTPUT_PATH = "../../NutriMatch_Presentation.pptx"

# Colors
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
MED_BLUE = RGBColor(0x4C, 0x72, 0xB0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
BLACK = RGBColor(0x00, 0x00, 0x00)
GREEN = RGBColor(0x55, 0xA8, 0x68)
RED = RGBColor(0xC4, 0x4E, 0x52)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_slide(title_text, layout_idx=6):
    """Add a blank slide and return it."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    return slide

def add_title_bar(slide, text, subtitle=None):
    """Add a colored title bar at the top."""
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)
        p2.alignment = PP_ALIGN.LEFT

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf

def add_bullet_slide(slide, bullets, left=Inches(0.8), top=Inches(1.5), width=Inches(11), height=Inches(5), font_size=20):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = BLACK
        p.space_after = Pt(8)
        p.level = 0
    return tf

def add_image_safe(slide, path, left, top, width=None, height=None):
    if os.path.isfile(path):
        if width and height:
            slide.shapes.add_picture(path, left, top, width, height)
        elif width:
            slide.shapes.add_picture(path, left, top, width=width)
        elif height:
            slide.shapes.add_picture(path, left, top, height=height)
        else:
            slide.shapes.add_picture(path, left, top)
        return True
    return False

def add_table(slide, data, left, top, col_widths, row_height=Inches(0.45)):
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top,
                                          sum(col_widths), row_height * rows)
    table = table_shape.table
    for i, w in enumerate(col_widths):
        table.columns[i].width = w
    for r_idx, row_data in enumerate(data):
        for c_idx, cell_text in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                if r_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.alignment = PP_ALIGN.CENTER
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_BLUE
            elif r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    return table

# ================================================================
# SLIDE 1: Title
# ================================================================
slide = add_slide("Title")
shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BLUE
shape.line.fill.background()

add_text_box(slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
             "NutriMatch", font_size=54, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
             "A Hybrid Recipe Recommendation System", font_size=28, color=RGBColor(0xBB,0xCC,0xDD), alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.8),
             "with Dietary Constraints and Explainable Recommendations", font_size=22, color=RGBColor(0x99,0xAA,0xBB), alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.6),
             "CS 513 Final Project", font_size=20, color=RGBColor(0x88,0x99,0xAA), alignment=PP_ALIGN.CENTER)

# ================================================================
# SLIDE 2: Problem Statement
# ================================================================
slide = add_slide("Problem")
add_title_bar(slide, "Problem Statement")
add_bullet_slide(slide, [
    "Recipe discovery is harder than generic search",
    "Users care about multiple constraints: taste, ingredients, time, nutrition, dietary restrictions",
    "A popular recipe may still be a poor recommendation for a specific user",
    "",
    "Research Question:",
    "How can a lightweight hybrid recommender combine user interaction history and recipe metadata",
    "to produce accurate, interpretable, and practically useful recipe suggestions?",
], font_size=20)

# ================================================================
# SLIDE 3: Dataset
# ================================================================
slide = add_slide("Dataset")
add_title_bar(slide, "Dataset: Food.com Recipes & Reviews")
data = [
    ["Metric", "Raw", "After Cleaning", "CF Subset"],
    ["Recipes", "231,637", "230,543", "39,844"],
    ["Interactions", "1,132,367", "1,067,281", "533,018"],
    ["Users", "-", "195,468", "16,973"],
    ["Sparsity", "-", "-", "99.92%"],
]
add_table(slide, data, Inches(1), Inches(1.8), [Inches(2.5), Inches(2.5), Inches(2.5), Inches(2.5)])

add_bullet_slide(slide, [
    "Features: ratings, ingredients, tags, nutrition, prep time, timestamps",
    "Temporal train/test split: each user's last interaction held out for testing",
    "CF subset: users and recipes with >= 5 ratings each",
], left=Inches(0.8), top=Inches(4.5), width=Inches(11), height=Inches(2.5), font_size=18)

# ================================================================
# SLIDE 4: EDA - Key Figures
# ================================================================
slide = add_slide("EDA")
add_title_bar(slide, "Exploratory Data Analysis")
add_image_safe(slide, os.path.join(A_FIGURES, "rating_distribution.png"),
               Inches(0.5), Inches(1.5), width=Inches(6), height=Inches(2.8))
add_image_safe(slide, os.path.join(A_FIGURES, "recipe_submission_trend.png"),
               Inches(6.8), Inches(1.5), width=Inches(6), height=Inches(2.8))
add_image_safe(slide, os.path.join(A_FIGURES, "interaction_density_distributions.png"),
               Inches(0.5), Inches(4.5), width=Inches(6), height=Inches(2.8))
add_image_safe(slide, os.path.join(A_FIGURES, "top_informative_tags.png"),
               Inches(6.8), Inches(4.5), width=Inches(6), height=Inches(2.8))

# ================================================================
# SLIDE 5: Method Overview
# ================================================================
slide = add_slide("Methods")
add_title_bar(slide, "Three Parallel Recommendation Approaches")

# Version A
add_text_box(slide, Inches(0.5), Inches(1.5), Inches(4), Inches(0.5),
             "Version A: Content-Based", font_size=22, bold=True, color=GREEN)
add_bullet_slide(slide, [
    "A0: Bayesian popularity baseline",
    "A2: TF-IDF item-to-item similarity",
    "A3: User-profile TF-IDF",
    "A4: Content + popularity reranker",
    "A5: SVD semantic content",
], left=Inches(0.5), top=Inches(2.1), width=Inches(3.8), height=Inches(3), font_size=16)

# Version B
add_text_box(slide, Inches(4.8), Inches(1.5), Inches(4), Inches(0.5),
             "Version B: Collaborative Filtering", font_size=22, bold=True, color=MED_BLUE)
add_bullet_slide(slide, [
    "B0: Bayesian popularity baseline",
    "B1: User-based kNN CF",
    "B2: Item-based kNN (implicit)",
    "B3: Item-based kNN (weighted)",
    "B4: SVD matrix factorization",
], left=Inches(4.8), top=Inches(2.1), width=Inches(3.8), height=Inches(3), font_size=16)

# Version C
add_text_box(slide, Inches(9.1), Inches(1.5), Inches(4), Inches(0.5),
             "Version C: Hybrid", font_size=22, bold=True, color=RED)
add_bullet_slide(slide, [
    "C0: SVD CF baseline",
    "C1: TF-IDF content baseline",
    "C2: Weighted hybrid (CF + content)",
    "C3: Switching hybrid",
    "C4: Reciprocal Rank Fusion",
], left=Inches(9.1), top=Inches(2.1), width=Inches(3.8), height=Inches(3), font_size=16)

add_text_box(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(1.5),
             "Evaluation: Same temporal split, same metrics (Precision@10, Recall@10, NDCG@10, Coverage@10) across all versions",
             font_size=18, color=DARK_BLUE)

# ================================================================
# SLIDE 6: Version A Results
# ================================================================
slide = add_slide("Version A")
add_title_bar(slide, "Version A: Content-Based Results")

a_data = [
    ["Model", "Precision@10", "Recall@10", "NDCG@10", "Coverage@10"],
    ["A0 Popularity", "0.00100", "0.01003", "0.00462", "0.06%"],
    ["A2 TF-IDF item", "0.00022", "0.00218", "0.00113", "44.85%"],
    ["A3 User-profile", "0.00027", "0.00268", "0.00132", "45.74%"],
    ["A4 Content+Pop", "0.00027", "0.00268", "0.00134", "45.30%"],
    ["A5 SVD semantic", "0.00021", "0.00209", "0.00096", "59.51%"],
]
add_table(slide, a_data, Inches(0.5), Inches(1.5),
          [Inches(2.5), Inches(2), Inches(2), Inches(2), Inches(2)])

add_image_safe(slide, os.path.join(A_FIGURES, "version_a_quality_vs_runtime_at_10.png"),
               Inches(0.5), Inches(4.3), width=Inches(6), height=Inches(3))

add_bullet_slide(slide, [
    "A4 selected as best Version A model",
    "Highest NDCG among personalized models",
    "45% catalog coverage (vs 0.06% for popularity)",
    "Strong for sparse users and cold-start items",
], left=Inches(7), top=Inches(4.3), width=Inches(5.5), height=Inches(3), font_size=16)

# ================================================================
# SLIDE 7: Version B Results
# ================================================================
slide = add_slide("Version B")
add_title_bar(slide, "Version B: Collaborative Filtering Results")

b_data = [
    ["Model", "Precision@10", "Recall@10", "NDCG@10", "Coverage@10"],
    ["B0 Popularity", "0.00103", "0.01027", "0.00471", "0.06%"],
    ["B1 User-kNN", "0.00193", "0.01927", "0.00984", "15.38%"],
    ["B2 Item-kNN impl", "0.00088", "0.00880", "0.00443", "46.77%"],
    ["B3 Item-kNN wgt", "0.00068", "0.00681", "0.00387", "54.50%"],
    ["B4 SVD (d=64)", "0.00252", "0.02524", "0.01309", "3.26%"],
]
add_table(slide, b_data, Inches(0.5), Inches(1.5),
          [Inches(2.5), Inches(2), Inches(2), Inches(2), Inches(2)])

add_bullet_slide(slide, [
    "B4 SVD is the strongest model in Version B",
    "NDCG@10 = 0.01309, ~10x better than best content model (A4)",
    "Matrix factorization outperforms all kNN variants",
    "Low coverage (3.26%) is the main weakness",
    "Trade-off: high precision but recommends from a narrow pool",
], left=Inches(0.5), top=Inches(4.5), width=Inches(12), height=Inches(2.5), font_size=18)

# ================================================================
# SLIDE 8: Version C Results
# ================================================================
slide = add_slide("Version C")
add_title_bar(slide, "Version C: Hybrid Recommendation Results")

c_data = [
    ["Model", "Precision@10", "Recall@10", "NDCG@10", "Coverage@10"],
    ["C0 SVD CF", "0.00252", "0.02524", "0.01309", "3.26%"],
    ["C1 TF-IDF content", "0.00011", "0.00105", "0.00053", "45.57%"],
    ["C2 Weighted hybrid", "0.00220", "0.02200", "0.01194", "17.31%"],
    ["C3 Switching hybrid", "0.00155", "0.01550", "0.00788", "40.63%"],
    ["C4 RRF hybrid", "0.00166", "0.01655", "0.00852", "29.12%"],
]
add_table(slide, c_data, Inches(0.5), Inches(1.5),
          [Inches(2.5), Inches(2), Inches(2), Inches(2), Inches(2)])

add_image_safe(slide, os.path.join(FIGURES_DIR, "version_c_ndcg_vs_coverage.png"),
               Inches(0.5), Inches(4.3), width=Inches(6), height=Inches(3))

add_bullet_slide(slide, [
    "C2 Weighted Hybrid is Version C's best model",
    "NDCG only 9% lower than pure CF",
    "Coverage 5x better: 17.3% vs 3.3%",
    "C3/C4 provide even higher coverage (29-41%)",
    "Hybrid methods successfully balance precision & diversity",
], left=Inches(7), top=Inches(4.3), width=Inches(5.5), height=Inches(3), font_size=16)

# ================================================================
# SLIDE 9: Version C Figures
# ================================================================
slide = add_slide("Version C Figures")
add_title_bar(slide, "Version C: Model Comparison Visualizations")

add_image_safe(slide, os.path.join(FIGURES_DIR, "version_c_model_metrics_at_10.png"),
               Inches(0.3), Inches(1.5), width=Inches(6.3), height=Inches(2.8))
add_image_safe(slide, os.path.join(FIGURES_DIR, "version_c_catalog_coverage_at_10.png"),
               Inches(6.8), Inches(1.5), width=Inches(6.2), height=Inches(2.8))
add_image_safe(slide, os.path.join(FIGURES_DIR, "version_c_hit_rate_at_10.png"),
               Inches(0.3), Inches(4.5), width=Inches(6.3), height=Inches(2.8))
add_image_safe(slide, os.path.join(FIGURES_DIR, "version_c_phase_runtime.png"),
               Inches(6.8), Inches(4.5), width=Inches(6.2), height=Inches(2.8))

# ================================================================
# SLIDE 10: Cross-Version Comparison
# ================================================================
slide = add_slide("Cross-Version")
add_title_bar(slide, "Cross-Version Comparison: A vs B vs C")

cross_data = [
    ["Version", "Best Model", "NDCG@10", "Recall@10", "Coverage@10", "Strength"],
    ["A (Content)", "A4 Reranker", "0.00134", "0.00268", "45.30%", "Cold-start, interpretable"],
    ["B (CF)", "B4 SVD", "0.01309", "0.02524", "3.26%", "Best precision"],
    ["C (Hybrid)", "C2 Weighted", "0.01194", "0.02200", "17.31%", "Balanced trade-off"],
]
add_table(slide, cross_data, Inches(0.5), Inches(1.5),
          [Inches(1.8), Inches(2), Inches(1.5), Inches(1.5), Inches(1.5), Inches(3)])

add_image_safe(slide, os.path.join(FIGURES_DIR, "cross_version_ndcg_vs_coverage.png"),
               Inches(0.5), Inches(3.8), width=Inches(6), height=Inches(3.5))
add_image_safe(slide, os.path.join(FIGURES_DIR, "cross_version_metrics_comparison.png"),
               Inches(6.8), Inches(3.8), width=Inches(6), height=Inches(3.5))

# ================================================================
# SLIDE 11: All Models Ranking
# ================================================================
slide = add_slide("All Models")
add_title_bar(slide, "All Models: NDCG@10 Ranking")
add_image_safe(slide, os.path.join(FIGURES_DIR, "all_models_ndcg_comparison.png"),
               Inches(1), Inches(1.5), width=Inches(11), height=Inches(5.5))

# ================================================================
# SLIDE 12: Key Findings
# ================================================================
slide = add_slide("Findings")
add_title_bar(slide, "Key Findings")
add_bullet_slide(slide, [
    "1. Collaborative filtering (SVD) significantly outperforms content-based methods",
    "     B4 SVD NDCG@10 is ~10x better than A4 content reranker",
    "",
    "2. Pure CF has very low catalog coverage (3.26%)",
    "     It only recommends from a small pool of popular items",
    "",
    "3. Hybrid methods successfully trade small precision loss for much better coverage",
    "     C2 weighted hybrid: only 9% NDCG drop, but 5x coverage improvement",
    "",
    "4. Switching and RRF hybrids provide even higher coverage (29-41%)",
    "     Better for diversity-focused applications",
    "",
    "5. Content-based methods are essential for cold-start and sparse users",
    "     C3 switching hybrid uses content for 43% of users who lack sufficient history",
], font_size=18)

# ================================================================
# SLIDE 13: Conclusion
# ================================================================
slide = add_slide("Conclusion")
add_title_bar(slide, "Conclusion & Recommendation")
add_bullet_slide(slide, [
    "Best overall system: C2 Weighted Hybrid (alpha=0.7)",
    "",
    "Why:",
    "    - Near-best NDCG (0.01194 vs 0.01309 for pure CF)",
    "    - 5x better catalog coverage (17.3% vs 3.3%)",
    "    - Lightweight and interpretable",
    "    - Combines the strengths of both content and CF approaches",
    "",
    "For production deployment:",
    "    - Use C2 as the default recommender",
    "    - Fall back to content-based (C1) for new/sparse users",
    "    - Use practical filters (prep time, dietary constraints) for better UX",
    "",
    "Future work:",
    "    - Deep learning models (e.g., neural collaborative filtering)",
    "    - LLM-based explanation layer for recommendation reasoning",
], font_size=18)

# ================================================================
# SLIDE 14: Thank You
# ================================================================
slide = add_slide("End")
shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BLUE
shape.line.fill.background()

add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
             "Thank You", font_size=54, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4), Inches(11), Inches(1),
             "Questions?", font_size=32, color=RGBColor(0xBB,0xCC,0xDD), alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.6),
             "CS 513 | NutriMatch | Hybrid Recipe Recommendation System", font_size=18,
             color=RGBColor(0x88,0x99,0xAA), alignment=PP_ALIGN.CENTER)

# ================================================================
# Save
# ================================================================
prs.save(OUTPUT_PATH)
print(f"Presentation saved to: {os.path.abspath(OUTPUT_PATH)}")
print(f"Total slides: {len(prs.slides)}")
