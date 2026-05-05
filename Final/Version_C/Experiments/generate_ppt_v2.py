"""Generate detailed presentation PPT for CS 513 Final Project - v2."""
import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

os.chdir(os.path.dirname(os.path.abspath(__file__)))

RESULTS_DIR = "../Results"
C_FIGURES = "../Results/Figures"
A_EDA_FIG = "../../Version_A/EDA/Figures"
A_RES_FIG = "../../Version_A/Results/Figures"
B_DIR = "../../Version_B"
OUTPUT_PATH = "../../NutriMatch_Presentation.pptx"

# Colors
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
BLACK = RGBColor(0x00, 0x00, 0x00)
GREEN_A = RGBColor(0x55, 0xA8, 0x68)
BLUE_B = RGBColor(0x4C, 0x72, 0xB0)
RED_C = RGBColor(0xC4, 0x4E, 0x52)
SUBTITLE_CLR = RGBColor(0xBB, 0xCC, 0xDD)
MUTED = RGBColor(0x88, 0x99, 0xAA)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def new_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])

def title_bar(slide, text, sub=None):
    s = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    s.fill.solid(); s.fill.fore_color.rgb = DARK_BLUE; s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.5); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(32); p.font.color.rgb = WHITE; p.font.bold = True
    if sub:
        p2 = tf.add_paragraph(); p2.text = sub; p2.font.size = Pt(18); p2.font.color.rgb = SUBTITLE_CLR

def txt(slide, l, t, w, h, text, sz=18, bold=False, clr=BLACK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz); p.font.bold = bold; p.font.color.rgb = clr; p.alignment = align
    return tf

def bullets(slide, items, l=Inches(0.8), t=Inches(1.5), w=Inches(11), h=Inches(5), sz=20):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(sz); p.font.color.rgb = BLACK; p.space_after = Pt(6)
    return tf

def img(slide, path, l, t, w=None, h=None):
    if not os.path.isfile(path): return False
    kw = {}
    if w: kw["width"] = w
    if h: kw["height"] = h
    slide.shapes.add_picture(path, l, t, **kw)
    return True

def table(slide, data, l, t, col_w, rh=Inches(0.42)):
    rows, cols = len(data), len(data[0])
    ts = slide.shapes.add_table(rows, cols, l, t, sum(col_w), rh * rows)
    tbl = ts.table
    for i, w in enumerate(col_w): tbl.columns[i].width = w
    for ri, rd in enumerate(data):
        for ci, ct in enumerate(rd):
            cell = tbl.cell(ri, ci); cell.text = str(ct)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13); p.alignment = PP_ALIGN.CENTER
                if ri == 0: p.font.bold = True; p.font.color.rgb = WHITE
            if ri == 0: cell.fill.solid(); cell.fill.fore_color.rgb = DARK_BLUE
            elif ri % 2 == 0: cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_GRAY

# ====================== SLIDE 1: Title ======================
sl = new_slide()
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg.fill.solid(); bg.fill.fore_color.rgb = DARK_BLUE; bg.line.fill.background()
txt(sl, Inches(1), Inches(1.8), Inches(11), Inches(1.5), "NutriMatch", 54, True, WHITE, PP_ALIGN.CENTER)
txt(sl, Inches(1), Inches(3.3), Inches(11), Inches(1), "A Hybrid Recipe Recommendation System", 28, False, SUBTITLE_CLR, PP_ALIGN.CENTER)
txt(sl, Inches(1), Inches(4.3), Inches(11), Inches(0.8), "with Dietary Constraints and Explainable Recommendations", 22, False, MUTED, PP_ALIGN.CENTER)
txt(sl, Inches(1), Inches(5.8), Inches(11), Inches(0.6), "CS 513 Final Project", 20, False, MUTED, PP_ALIGN.CENTER)

# ====================== SLIDE 2: Agenda ======================
sl = new_slide(); title_bar(sl, "Agenda")
bullets(sl, [
    "1.  Problem Statement & Motivation",
    "2.  Dataset Overview",
    "3.  Exploratory Data Analysis",
    "4.  Method Overview: Three Parallel Approaches",
    "5.  Version A: Content-Based Recommendation",
    "6.  Version B: Collaborative Filtering",
    "7.  Version C: Hybrid Recommendation",
    "8.  Cross-Version Comparison",
    "9.  Key Findings & Analysis",
    "10. Conclusion & Future Work",
], sz=22)

# ====================== SLIDE 3: Problem ======================
sl = new_slide(); title_bar(sl, "Problem Statement")
bullets(sl, [
    "Recipe discovery is harder than generic search:",
    "    Users care about taste, ingredients, time, nutrition, and dietary restrictions",
    "    A globally popular recipe may still be a poor fit for a specific user",
    "",
    "Research Questions:",
    "    1. How much does collaborative filtering improve over a popularity baseline?",
    "    2. How much do metadata features help in sparse-user or sparse-item settings?",
    "    3. Can a hybrid recommender balance precision, coverage, and practicality?",
    "",
    "Goal: Build a lightweight hybrid recommender that combines user history",
    "      and recipe metadata for accurate, diverse, and interpretable recommendations",
], sz=19)

# ====================== SLIDE 4: Dataset ======================
sl = new_slide(); title_bar(sl, "Dataset: Food.com Recipes & Reviews")
data = [
    ["Metric", "Raw Data", "After Cleaning", "CF Subset"],
    ["Recipes", "231,637", "230,543", "39,844"],
    ["Interactions", "1,132,367", "1,067,281", "533,018"],
    ["Users", "-", "195,468", "16,973"],
    ["Sparsity", "-", "-", "99.92%"],
    ["Train interactions", "-", "-", "507,043"],
    ["Test interactions", "-", "-", "10,069"],
]
table(sl, data, Inches(0.5), Inches(1.5), [Inches(2.5), Inches(2.5), Inches(2.5), Inches(2.5)])
bullets(sl, [
    "Rich metadata: ratings, ingredients, tags, nutrition, prep time, timestamps",
    "Temporal split: each user's last interaction held out for testing",
    "CF subset: users and recipes with >= 5 ratings each",
    "Two views: full metadata view (content) + filtered interaction view (CF)",
], l=Inches(0.5), t=Inches(5.2), w=Inches(12), h=Inches(2), sz=17)

# ====================== SLIDE 5: EDA page 1 ======================
sl = new_slide(); title_bar(sl, "EDA: Data Distributions")
img(sl, os.path.join(A_EDA_FIG, "rating_distribution.png"), Inches(0.3), Inches(1.5), w=Inches(6.3), h=Inches(2.7))
img(sl, os.path.join(A_EDA_FIG, "recipe_submission_trend.png"), Inches(6.8), Inches(1.5), w=Inches(6.2), h=Inches(2.7))
img(sl, os.path.join(A_EDA_FIG, "interaction_density_distributions.png"), Inches(0.3), Inches(4.4), w=Inches(6.3), h=Inches(2.7))
img(sl, os.path.join(A_EDA_FIG, "top_informative_tags.png"), Inches(6.8), Inches(4.4), w=Inches(6.2), h=Inches(2.7))

# ====================== SLIDE 6: EDA page 2 ======================
sl = new_slide(); title_bar(sl, "EDA: Metadata & Temporal Split")
img(sl, os.path.join(A_EDA_FIG, "recipe_metadata_distributions.png"), Inches(0.3), Inches(1.5), w=Inches(6.3), h=Inches(2.7))
img(sl, os.path.join(A_EDA_FIG, "correlation_heatmap.png"), Inches(6.8), Inches(1.5), w=Inches(6.2), h=Inches(2.7))
img(sl, os.path.join(A_EDA_FIG, "temporal_split_summary.png"), Inches(0.3), Inches(4.4), w=Inches(6.3), h=Inches(2.7))
img(sl, os.path.join(A_EDA_FIG, "history_bucket_overview.png"), Inches(6.8), Inches(4.4), w=Inches(6.2), h=Inches(2.7))

# ====================== SLIDE 7: EDA Sparsity (Version B) ======================
sl = new_slide(); title_bar(sl, "EDA: Data Sparsity for Collaborative Filtering")
img(sl, os.path.join(B_DIR, "bbd7a8feba91aa5b707f25ef43259d60.png"), Inches(1.5), Inches(1.5), w=Inches(5), h=Inches(3))
bullets(sl, [
    "16,973 users x 39,844 recipes",
    "533,018 interactions in the CF subset",
    "Density: 0.079% (extremely sparse)",
    "Sparsity: 99.92%",
    "",
    "This extreme sparsity is why:",
    "  - Pure CF struggles with coverage",
    "  - Content-based methods are needed for cold-start",
    "  - Hybrid approaches have potential value",
], l=Inches(7), t=Inches(1.5), w=Inches(5.5), h=Inches(5.5), sz=17)

# ====================== SLIDE 8: Method Overview ======================
sl = new_slide(); title_bar(sl, "Three Parallel Recommendation Approaches")
txt(sl, Inches(0.5), Inches(1.5), Inches(4), Inches(0.5), "Version A: Content-Based", 22, True, GREEN_A)
bullets(sl, [
    "A0: Bayesian popularity baseline",
    "A2: TF-IDF item-to-item similarity",
    "A3: User-profile TF-IDF",
    "A4: Content + popularity reranker",
    "A5: SVD semantic content",
], l=Inches(0.5), t=Inches(2.1), w=Inches(3.8), h=Inches(3), sz=15)

txt(sl, Inches(4.8), Inches(1.5), Inches(4), Inches(0.5), "Version B: Collaborative Filtering", 22, True, BLUE_B)
bullets(sl, [
    "B0: Bayesian popularity baseline",
    "B1: User-based kNN",
    "B2: Item-based kNN (implicit)",
    "B3: Item-based kNN (weighted)",
    "B4: SVD matrix factorization",
], l=Inches(4.8), t=Inches(2.1), w=Inches(3.8), h=Inches(3), sz=15)

txt(sl, Inches(9.1), Inches(1.5), Inches(4), Inches(0.5), "Version C: Hybrid", 22, True, RED_C)
bullets(sl, [
    "C0: SVD CF baseline",
    "C1: TF-IDF content baseline",
    "C2: Weighted hybrid (CF+content)",
    "C3: Switching hybrid",
    "C4: Reciprocal Rank Fusion",
], l=Inches(9.1), t=Inches(2.1), w=Inches(3.8), h=Inches(3), sz=15)

txt(sl, Inches(0.5), Inches(5.3), Inches(12), Inches(0.5),
    "Evaluation Protocol", 22, True, DARK_BLUE)
bullets(sl, [
    "Same temporal train/test split across all versions",
    "Metrics: Precision@10, Recall@10, NDCG@10, Catalog Coverage@10",
    "Positive relevance threshold: rating >= 4",
], l=Inches(0.5), t=Inches(5.9), w=Inches(12), h=Inches(1.5), sz=16)

# ====================== SLIDE 9: Version A Results ======================
sl = new_slide(); title_bar(sl, "Version A: Content-Based Results", "Lightweight metadata-first recommendation")
a_data = [
    ["Model", "Prec@10", "Recall@10", "NDCG@10", "Coverage@10"],
    ["A0 Popularity", "0.00100", "0.01003", "0.00462", "0.06%"],
    ["A2 TF-IDF item", "0.00022", "0.00218", "0.00113", "44.85%"],
    ["A3 User-profile", "0.00027", "0.00268", "0.00132", "45.74%"],
    ["A4 Content+Pop", "0.00027", "0.00268", "0.00134", "45.30%"],
    ["A5 SVD semantic", "0.00021", "0.00209", "0.00096", "59.51%"],
]
table(sl, a_data, Inches(0.5), Inches(1.5), [Inches(2.2), Inches(1.8), Inches(1.8), Inches(1.8), Inches(1.8)])
img(sl, os.path.join(A_RES_FIG, "version_a_quality_vs_runtime_at_10.png"), Inches(0.3), Inches(4.3), w=Inches(6.3), h=Inches(3))
bullets(sl, [
    "A4 selected as best Version A model",
    "Content models achieve 45-60% coverage",
    "But NDCG is low (~0.001)",
    "Strong for cold-start and sparse users",
    "Interpretable: can explain via ingredients/tags",
], l=Inches(7), t=Inches(4.3), w=Inches(5.5), h=Inches(3), sz=16)

# ====================== SLIDE 10: Version A Figures ======================
sl = new_slide(); title_bar(sl, "Version A: Detailed Visualizations")
img(sl, os.path.join(A_RES_FIG, "version_a_a0_popularity_tuning.png"), Inches(0.3), Inches(1.5), w=Inches(6.3), h=Inches(2.7))
img(sl, os.path.join(A_RES_FIG, "version_a_a3_tfidf_tuning.png"), Inches(6.8), Inches(1.5), w=Inches(6.2), h=Inches(2.7))
img(sl, os.path.join(A_RES_FIG, "version_a_content_explanation_terms.png"), Inches(0.3), Inches(4.4), w=Inches(6.3), h=Inches(2.7))
img(sl, os.path.join(A_RES_FIG, "version_a_hit_rank_distribution_at_10.png"), Inches(6.8), Inches(4.4), w=Inches(6.2), h=Inches(2.7))

# ====================== SLIDE 11: Version B Results ======================
sl = new_slide(); title_bar(sl, "Version B: Collaborative Filtering Results", "User interaction-based recommendation")
b_data = [
    ["Model", "Prec@10", "Recall@10", "NDCG@10", "Coverage@10"],
    ["B0 Popularity", "0.00103", "0.01027", "0.00471", "0.06%"],
    ["B1 User-kNN", "0.00193", "0.01927", "0.00984", "15.38%"],
    ["B2 Item-kNN impl", "0.00088", "0.00880", "0.00443", "46.77%"],
    ["B3 Item-kNN wgt", "0.00068", "0.00681", "0.00387", "54.50%"],
    ["B4 SVD (d=64)", "0.00252", "0.02524", "0.01309", "3.26%"],
]
table(sl, b_data, Inches(0.5), Inches(1.5), [Inches(2.2), Inches(1.8), Inches(1.8), Inches(1.8), Inches(1.8)])
bullets(sl, [
    "B4 SVD is the strongest model: NDCG ~10x better than A4",
    "Matrix factorization outperforms all kNN variants",
    "Trade-off: B4 has best precision but lowest coverage (3.26%)",
    "B2/B3 item-kNN have high coverage but lower precision",
], l=Inches(0.5), t=Inches(4.5), w=Inches(12), h=Inches(2.5), sz=18)

# ====================== SLIDE 12: Version B Tuning Figures ======================
sl = new_slide(); title_bar(sl, "Version B: Hyperparameter Tuning")
img(sl, os.path.join(B_DIR, "7ef17d3606b76bb0f9fc8ccd290cc0ba.png"), Inches(0.3), Inches(1.5), w=Inches(6.3), h=Inches(2.7))
img(sl, os.path.join(B_DIR, "994a688ad48fdc30fad7647e1f54396d.png"), Inches(6.8), Inches(1.5), w=Inches(6.2), h=Inches(2.7))
img(sl, os.path.join(B_DIR, "ac95ca5407f4da29b1b34656d7ab205d.png"), Inches(0.3), Inches(4.4), w=Inches(6.3), h=Inches(2.7))
bullets(sl, [
    "SVD: 64 components is optimal",
    "  128 overfits, 32 underfits",
    "",
    "kNN: User-kNN (B1) benefits from",
    "  more neighbors (k=50 best)",
    "  Item-kNN is stable across k",
    "",
    "Score distribution shows clear",
    "  ranking separation for top items",
], l=Inches(7), t=Inches(4.4), w=Inches(5.5), h=Inches(3), sz=15)

# ====================== SLIDE 13: Version C Results ======================
sl = new_slide(); title_bar(sl, "Version C: Hybrid Recommendation Results", "Combining CF and content signals")
c_data = [
    ["Model", "Prec@10", "Recall@10", "NDCG@10", "Coverage@10"],
    ["C0 SVD CF", "0.00252", "0.02524", "0.01309", "3.26%"],
    ["C1 TF-IDF content", "0.00011", "0.00105", "0.00053", "45.57%"],
    ["C2 Weighted hybrid", "0.00220", "0.02200", "0.01194", "17.31%"],
    ["C3 Switching hybrid", "0.00155", "0.01550", "0.00788", "40.63%"],
    ["C4 RRF hybrid", "0.00166", "0.01655", "0.00852", "29.12%"],
]
table(sl, c_data, Inches(0.5), Inches(1.5), [Inches(2.2), Inches(1.8), Inches(1.8), Inches(1.8), Inches(1.8)])
img(sl, os.path.join(C_FIGURES, "version_c_ndcg_vs_coverage.png"), Inches(0.3), Inches(4.3), w=Inches(6.3), h=Inches(3))
bullets(sl, [
    "C2 Weighted Hybrid is best trade-off:",
    "  Only 9% NDCG drop vs pure CF",
    "  5x better coverage (17% vs 3%)",
    "",
    "C3 Switching: 41% coverage, uses CF for",
    "  57% of users, content for rest",
    "",
    "C4 RRF: 29% coverage, rank-based fusion",
    "  robust to score scale differences",
], l=Inches(7), t=Inches(4.3), w=Inches(5.5), h=Inches(3), sz=15)

# ====================== SLIDE 14: Version C Figures ======================
sl = new_slide(); title_bar(sl, "Version C: Model Comparison Visualizations")
img(sl, os.path.join(C_FIGURES, "version_c_model_metrics_at_10.png"), Inches(0.3), Inches(1.5), w=Inches(6.3), h=Inches(2.7))
img(sl, os.path.join(C_FIGURES, "version_c_catalog_coverage_at_10.png"), Inches(6.8), Inches(1.5), w=Inches(6.2), h=Inches(2.7))
img(sl, os.path.join(C_FIGURES, "version_c_hit_rate_at_10.png"), Inches(0.3), Inches(4.4), w=Inches(6.3), h=Inches(2.7))
img(sl, os.path.join(C_FIGURES, "version_c_quality_vs_runtime_at_10.png"), Inches(6.8), Inches(4.4), w=Inches(6.2), h=Inches(2.7))

# ====================== SLIDE 15: Version C More Figures ======================
sl = new_slide(); title_bar(sl, "Version C: Runtime & Distribution Analysis")
img(sl, os.path.join(C_FIGURES, "version_c_phase_runtime.png"), Inches(0.3), Inches(1.5), w=Inches(6.3), h=Inches(2.7))
img(sl, os.path.join(C_FIGURES, "version_c_per_user_ndcg_boxplot.png"), Inches(6.8), Inches(1.5), w=Inches(6.2), h=Inches(2.7))
bullets(sl, [
    "C0 (pure SVD) is by far the fastest: 8 seconds total",
    "C1, C2, C4 each take ~10 minutes due to per-user cosine similarity over 230K recipes",
    "C3 switching is faster (4.5 min) because 57% of users skip content computation",
    "Per-user NDCG boxplot shows most users get 0 NDCG (extreme sparsity), but C0/C2 have higher upper-quartile hits",
], l=Inches(0.5), t=Inches(4.5), w=Inches(12), h=Inches(2.5), sz=17)

# ====================== SLIDE 16: Cross-Version Comparison ======================
sl = new_slide(); title_bar(sl, "Cross-Version Comparison: A vs B vs C")
cross_data = [
    ["Version", "Best Model", "NDCG@10", "Recall@10", "Coverage@10", "Strength"],
    ["A (Content)", "A4 Reranker", "0.00134", "0.00268", "45.30%", "Cold-start, interpretable"],
    ["B (CF)", "B4 SVD", "0.01309", "0.02524", "3.26%", "Best precision"],
    ["C (Hybrid)", "C2 Weighted", "0.01194", "0.02200", "17.31%", "Balanced trade-off"],
]
table(sl, cross_data, Inches(0.5), Inches(1.5), [Inches(1.8), Inches(2), Inches(1.5), Inches(1.5), Inches(1.5), Inches(3)])
img(sl, os.path.join(C_FIGURES, "cross_version_metrics_comparison.png"), Inches(0.3), Inches(3.8), w=Inches(6.3), h=Inches(3.5))
img(sl, os.path.join(C_FIGURES, "cross_version_ndcg_vs_coverage.png"), Inches(6.8), Inches(3.8), w=Inches(6.2), h=Inches(3.5))

# ====================== SLIDE 17: Cross-Version Coverage + All Models ======================
sl = new_slide(); title_bar(sl, "Cross-Version: Coverage & All Models Ranking")
img(sl, os.path.join(C_FIGURES, "cross_version_coverage_comparison.png"), Inches(0.3), Inches(1.5), w=Inches(6.3), h=Inches(3))
img(sl, os.path.join(C_FIGURES, "all_models_ndcg_comparison.png"), Inches(6.8), Inches(1.5), w=Inches(6.2), h=Inches(5.5))
bullets(sl, [
    "Content methods (A): highest coverage, lowest NDCG",
    "CF methods (B): highest NDCG, lowest coverage",
    "Hybrid (C2): best balance between the two extremes",
], l=Inches(0.3), t=Inches(4.8), w=Inches(6.3), h=Inches(2.5), sz=16)

# ====================== SLIDE 18: Key Findings ======================
sl = new_slide(); title_bar(sl, "Key Findings")
bullets(sl, [
    "1.  Collaborative filtering (SVD) significantly outperforms content-based methods",
    "       B4 SVD NDCG@10 = 0.01309, roughly 10x better than best content model A4 (0.00134)",
    "",
    "2.  Pure CF suffers from extremely low catalog coverage (3.26%)",
    "       It recommends from a very narrow pool of popular items",
    "",
    "3.  Hybrid methods successfully trade small precision loss for much better coverage",
    "       C2 weighted hybrid: only 9% NDCG drop, but 5x coverage improvement (17.3%)",
    "",
    "4.  Switching and RRF hybrids achieve even higher coverage (29-41%)",
    "       C3 uses content for 43% of users who lack sufficient interaction history",
    "",
    "5.  The precision-coverage trade-off is the central design decision",
    "       Application context determines which model to deploy",
], sz=18)

# ====================== SLIDE 19: Practical Implications ======================
sl = new_slide(); title_bar(sl, "Practical Implications")
bullets(sl, [
    "For a real recipe recommendation app:",
    "",
    "    High-precision scenario (e.g., homepage 'Top Picks for You'):",
    "        Use B4/C0 (SVD) for users with sufficient history",
    "",
    "    Discovery scenario (e.g., 'Explore New Recipes'):",
    "        Use C3/C4 hybrid for better diversity and coverage",
    "",
    "    Cold-start users (new sign-ups):",
    "        Use C1/A4 content-based until enough ratings are collected",
    "",
    "    With dietary constraints:",
    "        Apply post-filters for prep time, ingredients, nutrition",
    "        Content-based methods naturally support these constraints",
], sz=18)

# ====================== SLIDE 20: Conclusion ======================
sl = new_slide(); title_bar(sl, "Conclusion & Future Work")
bullets(sl, [
    "Conclusion:",
    "    Best overall system: C2 Weighted Hybrid (alpha=0.7)",
    "    Near-best NDCG (0.01194 vs 0.01309 for pure CF)",
    "    5x better catalog coverage (17.3% vs 3.3%)",
    "    Lightweight, interpretable, and practical",
    "",
    "Limitations:",
    "    All models show low absolute metrics due to extreme data sparsity (99.92%)",
    "    Only one temporal split evaluated; cross-validation would strengthen results",
    "    No deep learning models compared (e.g., neural CF, transformers)",
    "",
    "Future Work:",
    "    Neural collaborative filtering (NCF) for better latent representations",
    "    LLM-based explanation layer for recommendation reasoning",
    "    Online A/B testing to validate offline metrics in real user scenarios",
    "    Incorporate side information: images, cooking steps, user demographics",
], sz=17)

# ====================== SLIDE 21: Thank You ======================
sl = new_slide()
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg.fill.solid(); bg.fill.fore_color.rgb = DARK_BLUE; bg.line.fill.background()
txt(sl, Inches(1), Inches(2.2), Inches(11), Inches(1.5), "Thank You", 54, True, WHITE, PP_ALIGN.CENTER)
txt(sl, Inches(1), Inches(3.8), Inches(11), Inches(1), "Questions?", 32, False, SUBTITLE_CLR, PP_ALIGN.CENTER)
txt(sl, Inches(1), Inches(5.5), Inches(11), Inches(0.6),
    "CS 513  |  NutriMatch  |  Hybrid Recipe Recommendation System", 18, False, MUTED, PP_ALIGN.CENTER)

# ====================== SAVE ======================
prs.save(OUTPUT_PATH)
print(f"Saved: {os.path.abspath(OUTPUT_PATH)}")
print(f"Total slides: {len(prs.slides)}")
