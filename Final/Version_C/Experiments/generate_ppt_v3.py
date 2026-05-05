"""Generate comprehensive presentation PPT for CS 513 Final Project - v3.
All figures included, with coherent narrative flow."""
import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

os.chdir(os.path.dirname(os.path.abspath(__file__)))

RESULTS_DIR = "../Results"
C_FIG = "../Results/Figures"
A_EDA = "../../Version_A/EDA/Figures"
A_FIG = "../../Version_A/Results/Figures"
B_DIR = "../../Version_B"
OUTPUT_PATH = "../../NutriMatch_Presentation.pptx"

DARK = RGBColor(0x1B, 0x3A, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xF0, 0xF0, 0xF0)
BLK = RGBColor(0x00, 0x00, 0x00)
GA = RGBColor(0x55, 0xA8, 0x68)
BB = RGBColor(0x4C, 0x72, 0xB0)
RC = RGBColor(0xC4, 0x4E, 0x52)
SUB = RGBColor(0xBB, 0xCC, 0xDD)
MUT = RGBColor(0x88, 0x99, 0xAA)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

def ns(): return prs.slides.add_slide(prs.slide_layouts[6])

def tb(sl, text, sub=None):
    s = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    s.fill.solid(); s.fill.fore_color.rgb = DARK; s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.5); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(32); p.font.color.rgb = WHITE; p.font.bold = True
    if sub:
        p2 = tf.add_paragraph(); p2.text = sub; p2.font.size = Pt(18); p2.font.color.rgb = SUB

def tx(sl, l, t, w, h, text, sz=18, b=False, c=BLK, a=PP_ALIGN.LEFT):
    bx = sl.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz); p.font.bold = b; p.font.color.rgb = c; p.alignment = a
    return tf

def bl(sl, items, l=Inches(0.8), t=Inches(1.5), w=Inches(11), h=Inches(5), sz=20):
    bx = sl.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(sz); p.font.color.rgb = BLK; p.space_after = Pt(6)
    return tf

def im(sl, path, l, t, w=None, h=None):
    if not os.path.isfile(path): print(f"  MISSING: {path}"); return False
    kw = {};
    if w: kw["width"] = w
    if h: kw["height"] = h
    sl.shapes.add_picture(path, l, t, **kw); return True

def tbl(sl, data, l, t, cw, rh=Inches(0.42)):
    rows, cols = len(data), len(data[0])
    ts = sl.shapes.add_table(rows, cols, l, t, sum(cw), rh * rows)
    tb = ts.table
    for i, w in enumerate(cw): tb.columns[i].width = w
    for ri, rd in enumerate(data):
        for ci, ct in enumerate(rd):
            cell = tb.cell(ri, ci); cell.text = str(ct)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13); p.alignment = PP_ALIGN.CENTER
                if ri == 0: p.font.bold = True; p.font.color.rgb = WHITE
            if ri == 0: cell.fill.solid(); cell.fill.fore_color.rgb = DARK
            elif ri % 2 == 0: cell.fill.solid(); cell.fill.fore_color.rgb = LGRAY

# ======================== 1. TITLE ========================
sl = ns()
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg.fill.solid(); bg.fill.fore_color.rgb = DARK; bg.line.fill.background()
tx(sl, Inches(1), Inches(1.8), Inches(11), Inches(1.5), "NutriMatch", 54, True, WHITE, PP_ALIGN.CENTER)
tx(sl, Inches(1), Inches(3.3), Inches(11), Inches(1), "A Hybrid Recipe Recommendation System", 28, False, SUB, PP_ALIGN.CENTER)
tx(sl, Inches(1), Inches(4.3), Inches(11), Inches(0.8), "with Dietary Constraints and Explainable Recommendations", 22, False, MUT, PP_ALIGN.CENTER)
tx(sl, Inches(1), Inches(5.8), Inches(11), Inches(0.6), "CS 513 Final Project", 20, False, MUT, PP_ALIGN.CENTER)

# ======================== 2. AGENDA ========================
sl = ns(); tb(sl, "Agenda")
bl(sl, [
    "1.   Problem & Motivation",
    "2.   Dataset & Preprocessing",
    "3.   Exploratory Data Analysis",
    "4.   Methodology: Three Parallel Approaches",
    "5.   Version A: Content-Based Recommendation",
    "6.   Version B: Collaborative Filtering",
    "7.   Version C: Hybrid Recommendation",
    "8.   Cross-Version Comparison & Analysis",
    "9.   Key Findings",
    "10.  Practical Implications",
    "11.  Conclusion & Future Work",
], sz=22)

# ======================== 3. PROBLEM ========================
sl = ns(); tb(sl, "Problem Statement & Motivation")
bl(sl, [
    "The Challenge:",
    "    Recipe discovery is harder than generic search because users have multiple constraints:",
    "    taste preferences, available ingredients, cooking time, nutrition goals, dietary restrictions",
    "",
    "Why existing approaches fall short:",
    "    Popularity-based: ignores personal preferences entirely",
    "    Content-based: can't capture 'people like you also liked...' patterns",
    "    Collaborative filtering: fails for new users or items with no history (cold start)",
    "",
    "Our Research Questions:",
    "    Q1. How much does collaborative filtering improve over a popularity baseline?",
    "    Q2. How much do metadata features help for sparse users and items?",
    "    Q3. Can a hybrid recommender balance precision, coverage, and practicality?",
], sz=18)

# ======================== 4. DATASET ========================
sl = ns(); tb(sl, "Dataset: Food.com Recipes & Reviews")
data = [
    ["Metric", "Raw Data", "After Cleaning", "CF Subset"],
    ["Recipes", "231,637", "230,543", "39,844"],
    ["Interactions", "1,132,367", "1,067,281", "533,018"],
    ["Users", "-", "195,468", "16,973"],
    ["Sparsity", "-", "-", "99.92%"],
    ["Train set", "-", "-", "507,043"],
    ["Test set", "-", "-", "10,069"],
]
tbl(sl, data, Inches(0.5), Inches(1.5), [Inches(2.5), Inches(2.5), Inches(2.5), Inches(2.5)])
bl(sl, [
    "Features: explicit ratings, ingredients, tags, nutrition, prep time, timestamps",
    "Temporal split: each user's last interaction held out for evaluation",
    "CF subset: users and recipes with >= 5 ratings each",
    "Two data views: full metadata (content) + filtered interactions (CF)",
    "This allows honest study of cold-start and sparsity, not just one preprocessing choice",
], l=Inches(0.5), t=Inches(5.0), w=Inches(12), h=Inches(2.3), sz=17)

# ======================== 5. EDA: Ratings & Trends ========================
sl = ns(); tb(sl, "EDA: Rating Distribution & Submission Trends",
              "Understanding the shape of the data before modeling")
im(sl, os.path.join(A_EDA, "rating_distribution.png"), Inches(0.3), Inches(1.5), w=Inches(6.3), h=Inches(2.7))
im(sl, os.path.join(A_EDA, "recipe_submission_trend.png"), Inches(6.8), Inches(1.5), w=Inches(6.2), h=Inches(2.7))
bl(sl, [
    "Left: Ratings are heavily skewed toward 5 stars - most interactions are positive",
    "    This means binary positive/negative signal is more useful than fine-grained rating prediction",
    "Right: Recipe submissions peaked around 2008-2010, then declined",
    "    Older recipes dominate the catalog; newer recipes may have less history",
], l=Inches(0.5), t=Inches(4.5), w=Inches(12), h=Inches(2.5), sz=16)

# ======================== 6. EDA: Interaction Density & Tags ========================
sl = ns(); tb(sl, "EDA: Interaction Density & Content Features",
              "Identifying sparsity patterns and informative metadata")
im(sl, os.path.join(A_EDA, "interaction_density_distributions.png"), Inches(0.3), Inches(1.5), w=Inches(6.3), h=Inches(2.7))
im(sl, os.path.join(A_EDA, "top_informative_tags.png"), Inches(6.8), Inches(1.5), w=Inches(6.2), h=Inches(2.7))
bl(sl, [
    "Left: Long-tail distribution - most users and recipes have very few interactions",
    "    This justifies the need for content features to supplement sparse interaction data",
    "Right: Tags like cuisine type, dietary labels, and meal type are highly informative",
    "    These features power content-based and hybrid recommendation",
], l=Inches(0.5), t=Inches(4.5), w=Inches(12), h=Inches(2.5), sz=16)

# ======================== 7. EDA: Metadata & Correlation ========================
sl = ns(); tb(sl, "EDA: Recipe Metadata & Feature Correlations",
              "Examining the richness of recipe-level features")
im(sl, os.path.join(A_EDA, "recipe_metadata_distributions.png"), Inches(0.3), Inches(1.5), w=Inches(6.3), h=Inches(2.7))
im(sl, os.path.join(A_EDA, "correlation_heatmap.png"), Inches(6.8), Inches(1.5), w=Inches(6.2), h=Inches(2.7))
bl(sl, [
    "Left: Prep time, ingredient count, and step count vary widely across recipes",
    "    These features support practical filters (quick meals, simple recipes)",
    "Right: Correlation heatmap shows nutrition features are moderately correlated",
    "    Ingredient count and step count correlate with complexity metrics",
], l=Inches(0.5), t=Inches(4.5), w=Inches(12), h=Inches(2.5), sz=16)

# ======================== 8. EDA: Temporal Split & History ========================
sl = ns(); tb(sl, "EDA: Temporal Split & User History Buckets",
              "Setting up a realistic evaluation framework")
im(sl, os.path.join(A_EDA, "temporal_split_summary.png"), Inches(0.3), Inches(1.5), w=Inches(6.3), h=Inches(2.7))
im(sl, os.path.join(A_EDA, "history_bucket_overview.png"), Inches(6.8), Inches(1.5), w=Inches(6.2), h=Inches(2.7))
bl(sl, [
    "Left: Per-user temporal holdout - each user's last interaction is the test item",
    "    This respects chronological order, unlike random splits that cause data leakage",
    "Right: Users grouped by history size - many users have minimal history",
    "    This motivates hybrid methods that handle sparse and rich users differently",
], l=Inches(0.5), t=Inches(4.5), w=Inches(12), h=Inches(2.5), sz=16)

# ======================== 9. EDA: Sparsity Context ========================
sl = ns(); tb(sl, "EDA: Data Sparsity for Collaborative Filtering",
              "Why sparsity is the fundamental challenge")
im(sl, os.path.join(B_DIR, "bbd7a8feba91aa5b707f25ef43259d60.png"), Inches(1.5), Inches(1.5), w=Inches(5), h=Inches(3))
bl(sl, [
    "16,973 users x 39,844 recipes",
    "533,018 interactions out of 676M possible",
    "Density: 0.079%  |  Sparsity: 99.92%",
    "",
    "Implications for modeling:",
    "  Pure CF can only work where history exists",
    "  Content features are needed for the 99.9% empty cells",
    "  This is the core motivation for hybrid approaches",
], l=Inches(7), t=Inches(1.5), w=Inches(5.5), h=Inches(5.5), sz=17)

# ======================== 10. METHOD OVERVIEW ========================
sl = ns(); tb(sl, "Methodology: Three Parallel Approaches",
              "Shared data foundation + independent modeling tracks")
tx(sl, Inches(0.5), Inches(1.5), Inches(4), Inches(0.5), "Version A: Content-Based", 22, True, GA)
bl(sl, ["A0: Bayesian popularity", "A2: TF-IDF item-to-item", "A3: User-profile TF-IDF",
        "A4: Content + popularity reranker", "A5: SVD semantic content"],
   l=Inches(0.5), t=Inches(2.1), w=Inches(3.8), h=Inches(2.5), sz=15)

tx(sl, Inches(4.8), Inches(1.5), Inches(4), Inches(0.5), "Version B: Collaborative Filtering", 22, True, BB)
bl(sl, ["B0: Bayesian popularity", "B1: User-based kNN", "B2: Item-based kNN (implicit)",
        "B3: Item-based kNN (weighted)", "B4: SVD matrix factorization"],
   l=Inches(4.8), t=Inches(2.1), w=Inches(3.8), h=Inches(2.5), sz=15)

tx(sl, Inches(9.1), Inches(1.5), Inches(4), Inches(0.5), "Version C: Hybrid", 22, True, RC)
bl(sl, ["C0: SVD CF baseline", "C1: TF-IDF content baseline", "C2: Weighted hybrid",
        "C3: Switching hybrid", "C4: Reciprocal Rank Fusion"],
   l=Inches(9.1), t=Inches(2.1), w=Inches(3.8), h=Inches(2.5), sz=15)

tx(sl, Inches(0.5), Inches(4.8), Inches(12), Inches(0.5), "Shared Evaluation Protocol", 22, True, DARK)
bl(sl, [
    "Same temporal train/test split  |  Same positive threshold (rating >= 4)",
    "Metrics: Precision@10, Recall@10, NDCG@10, Catalog Coverage@10, Hit Rate@10",
    "All models exclude training items from recommendations (no leakage)",
], l=Inches(0.5), t=Inches(5.4), w=Inches(12), h=Inches(2), sz=16)

# ======================== 11. VERSION A: Data Snapshot ========================
sl = ns(); tb(sl, "Version A: Data Foundation",
              "Understanding the data before building content models")
im(sl, os.path.join(A_FIG, "version_a_primary_data_snapshot.png"), Inches(0.3), Inches(1.5), w=Inches(6.3), h=Inches(2.7))
im(sl, os.path.join(A_FIG, "version_a_train_test_rating_distribution.png"), Inches(6.8), Inches(1.5), w=Inches(6.2), h=Inches(2.7))
im(sl, os.path.join(A_FIG, "version_a_eval_user_selection.png"), Inches(0.3), Inches(4.4), w=Inches(6.3), h=Inches(2.7))
im(sl, os.path.join(A_FIG, "version_a_content_text_length_distribution.png"), Inches(6.8), Inches(4.4), w=Inches(6.2), h=Inches(2.7))

# ======================== 12. VERSION A: Results Table ========================
sl = ns(); tb(sl, "Version A: Content-Based Results",
              "5 models evaluated on the filtered temporal split")
a_data = [
    ["Model", "Prec@10", "Recall@10", "NDCG@10", "Coverage@10", "Sec/User"],
    ["A0 Bayesian Pop", "0.00100", "0.01003", "0.00462", "0.06%", "0.0035"],
    ["A2 TF-IDF item", "0.00022", "0.00218", "0.00113", "44.85%", "0.286"],
    ["A3 User-profile", "0.00027", "0.00268", "0.00132", "45.74%", "0.058"],
    ["A4 Content+Pop", "0.00027", "0.00268", "0.00134", "45.30%", "0.061"],
    ["A5 SVD semantic", "0.00021", "0.00209", "0.00096", "59.51%", "0.016"],
]
tbl(sl, a_data, Inches(0.3), Inches(1.5), [Inches(2.2), Inches(1.7), Inches(1.7), Inches(1.7), Inches(1.7), Inches(1.7)])
im(sl, os.path.join(A_FIG, "version_a_model_metrics_at_10.png"), Inches(0.3), Inches(4.3), w=Inches(6.3), h=Inches(3))
bl(sl, [
    "A4 selected as Version A's best model:",
    "  Best NDCG among personalized models",
    "  45% coverage (vs 0.06% for popularity)",
    "  Lightweight: 0.06 sec/user",
    "",
    "Key insight: content models trade precision",
    "for much broader item coverage",
], l=Inches(7), t=Inches(4.3), w=Inches(5.5), h=Inches(3), sz=15)

# ======================== 13. VERSION A: Tuning ========================
sl = ns(); tb(sl, "Version A: Hyperparameter Tuning",
              "Finding optimal parameters for each content model")
im(sl, os.path.join(A_FIG, "version_a_a0_popularity_tuning.png"), Inches(0.3), Inches(1.5), w=Inches(6.3), h=Inches(2.7))
im(sl, os.path.join(A_FIG, "version_a_a3_tfidf_tuning.png"), Inches(6.8), Inches(1.5), w=Inches(6.2), h=Inches(2.7))
im(sl, os.path.join(A_FIG, "version_a_a4_rerank_weight_tuning.png"), Inches(0.3), Inches(4.4), w=Inches(6.3), h=Inches(2.7))
im(sl, os.path.join(A_FIG, "version_a_a4_rerank_weight_heatmap.png"), Inches(6.8), Inches(4.4), w=Inches(6.2), h=Inches(2.7))

# ======================== 14. VERSION A: Analysis ========================
sl = ns(); tb(sl, "Version A: Quality, Runtime & Explainability",
              "Content models offer interpretability that CF cannot")
im(sl, os.path.join(A_FIG, "version_a_quality_vs_runtime_at_10.png"), Inches(0.3), Inches(1.5), w=Inches(6.3), h=Inches(2.7))
im(sl, os.path.join(A_FIG, "version_a_phase_runtime.png"), Inches(6.8), Inches(1.5), w=Inches(6.2), h=Inches(2.7))
im(sl, os.path.join(A_FIG, "version_a_content_explanation_terms.png"), Inches(0.3), Inches(4.4), w=Inches(6.3), h=Inches(2.7))
im(sl, os.path.join(A_FIG, "version_a_constraint_demo_minutes.png"), Inches(6.8), Inches(4.4), w=Inches(6.2), h=Inches(2.7))

# ======================== 15. VERSION A: Hit Analysis ========================
sl = ns(); tb(sl, "Version A: Hit Distribution & Recommendation Examples",
              "Where content models succeed and where they struggle")
im(sl, os.path.join(A_FIG, "version_a_hit_rank_distribution_at_10.png"), Inches(0.3), Inches(1.5), w=Inches(6.3), h=Inches(2.7))
im(sl, os.path.join(A_FIG, "version_a_sample_recommendation_scores.png"), Inches(6.8), Inches(1.5), w=Inches(6.2), h=Inches(2.7))
im(sl, os.path.join(A_FIG, "version_a_top_bayesian_popularity_recipes.png"), Inches(0.3), Inches(4.4), w=Inches(6.3), h=Inches(2.7))
bl(sl, [
    "Hit rank distribution: A0 (popularity)",
    "concentrates hits at rank 3 - when it",
    "hits, it hits early",
    "",
    "Content models hit more uniformly",
    "across all 10 positions",
    "",
    "Popularity baseline recommends the",
    "same few recipes to everyone - hence",
    "the 0.06% coverage",
], l=Inches(7), t=Inches(4.4), w=Inches(5.5), h=Inches(2.7), sz=14)

# ======================== 16. VERSION A: Summary ========================
sl = ns(); tb(sl, "Version A: Summary & Transition to Version B")
bl(sl, [
    "Version A established the content-based baseline:",
    "",
    "    Strengths:",
    "        High catalog coverage (45-60%) - recommends diverse recipes",
    "        Interpretable - can explain via shared ingredients, tags, nutrition",
    "        Handles cold-start users and items naturally",
    "        Supports practical filters (prep time, dietary constraints)",
    "",
    "    Weaknesses:",
    "        Low absolute NDCG (best = 0.00134) - not very accurate",
    "        Cannot capture collaborative signals ('people like you also liked...')",
    "",
    "    Question for Version B:",
    "        How much does adding user interaction patterns improve precision?",
    "        And what do we lose in return?",
], sz=18)

# ======================== 17. VERSION B: Results ========================
sl = ns(); tb(sl, "Version B: Collaborative Filtering Results",
              "5 models using user-item interaction patterns")
b_data = [
    ["Model", "Prec@10", "Recall@10", "NDCG@10", "Coverage@10", "Sec/User"],
    ["B0 Popularity", "0.00103", "0.01027", "0.00471", "0.06%", "0.000005"],
    ["B1 User-kNN", "0.00193", "0.01927", "0.00984", "15.38%", "0.000614"],
    ["B2 Item-kNN impl", "0.00088", "0.00880", "0.00443", "46.77%", "0.003000"],
    ["B3 Item-kNN wgt", "0.00068", "0.00681", "0.00387", "54.50%", "0.001521"],
    ["B4 SVD (d=64)", "0.00252", "0.02524", "0.01309", "3.26%", "0.000896"],
]
tbl(sl, b_data, Inches(0.3), Inches(1.5), [Inches(2.2), Inches(1.7), Inches(1.7), Inches(1.7), Inches(1.7), Inches(1.7)])
bl(sl, [
    "B4 SVD is the clear winner: NDCG = 0.01309, roughly 10x better than A4",
    "Matrix factorization captures latent user-item patterns that content features miss",
    "But B4 coverage is only 3.26% - it recommends from a very narrow pool",
    "Item-kNN models (B2/B3) have high coverage but poor precision",
    "Answer to Q1: CF improves dramatically over the popularity baseline",
], l=Inches(0.5), t=Inches(4.5), w=Inches(12), h=Inches(2.5), sz=17)

# ======================== 18. VERSION B: Tuning ========================
sl = ns(); tb(sl, "Version B: Hyperparameter Tuning & Analysis",
              "SVD components and kNN neighbor count optimization")
im(sl, os.path.join(B_DIR, "7ef17d3606b76bb0f9fc8ccd290cc0ba.png"), Inches(0.3), Inches(1.5), w=Inches(6.3), h=Inches(2.7))
im(sl, os.path.join(B_DIR, "994a688ad48fdc30fad7647e1f54396d.png"), Inches(6.8), Inches(1.5), w=Inches(6.2), h=Inches(2.7))
im(sl, os.path.join(B_DIR, "ac95ca5407f4da29b1b34656d7ab205d.png"), Inches(0.3), Inches(4.4), w=Inches(6.3), h=Inches(2.7))
bl(sl, [
    "SVD tuning (top-left):",
    "  64 components is optimal",
    "  128 overfits, 32 underfits",
    "",
    "kNN tuning (top-right):",
    "  User-kNN improves with more",
    "  neighbors; item-kNN is stable",
    "",
    "Score distribution (bottom-left):",
    "  Clear ranking separation between",
    "  top-ranked and lower items",
], l=Inches(7), t=Inches(4.4), w=Inches(5.5), h=Inches(2.7), sz=14)

# ======================== 19. VERSION B: Summary ========================
sl = ns(); tb(sl, "Version B: Summary & Transition to Version C")
bl(sl, [
    "Version B demonstrated the power of collaborative filtering:",
    "",
    "    Strengths:",
    "        B4 SVD achieves NDCG 10x higher than best content model",
    "        Fast inference: < 1ms per user recommendation",
    "        Captures real user preference patterns from interaction history",
    "",
    "    Weaknesses:",
    "        Very low catalog coverage (3.26%) - echo chamber risk",
    "        Cannot handle cold-start users or items with no history",
    "        Not interpretable - can't explain why a recipe was recommended",
    "",
    "    The natural next question:",
    "        Can we combine CF's precision with content's coverage?",
    "        This is exactly what Version C attempts with hybrid methods",
], sz=18)

# ======================== 20. VERSION C: Results ========================
sl = ns(); tb(sl, "Version C: Hybrid Recommendation Results",
              "Combining the strengths of content and collaborative filtering")
c_data = [
    ["Model", "Prec@10", "Recall@10", "NDCG@10", "Coverage@10", "Method"],
    ["C0 SVD CF", "0.00252", "0.02524", "0.01309", "3.26%", "Pure CF (baseline)"],
    ["C1 TF-IDF", "0.00011", "0.00105", "0.00053", "45.57%", "Pure content (baseline)"],
    ["C2 Weighted", "0.00220", "0.02200", "0.01194", "17.31%", "0.7*CF + 0.3*content"],
    ["C3 Switching", "0.00155", "0.01550", "0.00788", "40.63%", "CF or content per user"],
    ["C4 RRF", "0.00166", "0.01655", "0.00852", "29.12%", "Rank fusion"],
]
tbl(sl, c_data, Inches(0.3), Inches(1.5), [Inches(1.8), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5), Inches(3)])
im(sl, os.path.join(C_FIG, "version_c_ndcg_vs_coverage.png"), Inches(0.3), Inches(4.3), w=Inches(6.3), h=Inches(3))
bl(sl, [
    "C2 Weighted Hybrid is the best trade-off:",
    "  NDCG = 0.01194 (only 9% below pure CF)",
    "  Coverage = 17.3% (5x better than pure CF)",
    "",
    "C3 Switching: 41% coverage, adapts per user",
    "  CF for 57% of users, content for 43%",
    "",
    "C4 RRF: 29% coverage, score-scale agnostic",
], l=Inches(7), t=Inches(4.3), w=Inches(5.5), h=Inches(3), sz=15)

# ======================== 21. VERSION C: Figures 1 ========================
sl = ns(); tb(sl, "Version C: Model Comparison Visualizations")
im(sl, os.path.join(C_FIG, "version_c_model_metrics_at_10.png"), Inches(0.3), Inches(1.5), w=Inches(6.3), h=Inches(2.7))
im(sl, os.path.join(C_FIG, "version_c_catalog_coverage_at_10.png"), Inches(6.8), Inches(1.5), w=Inches(6.2), h=Inches(2.7))
im(sl, os.path.join(C_FIG, "version_c_hit_rate_at_10.png"), Inches(0.3), Inches(4.4), w=Inches(6.3), h=Inches(2.7))
im(sl, os.path.join(C_FIG, "version_c_quality_vs_runtime_at_10.png"), Inches(6.8), Inches(4.4), w=Inches(6.2), h=Inches(2.7))

# ======================== 22. VERSION C: Figures 2 ========================
sl = ns(); tb(sl, "Version C: Runtime & Per-User Distribution")
im(sl, os.path.join(C_FIG, "version_c_phase_runtime.png"), Inches(0.3), Inches(1.5), w=Inches(6.3), h=Inches(2.7))
im(sl, os.path.join(C_FIG, "version_c_per_user_ndcg_boxplot.png"), Inches(6.8), Inches(1.5), w=Inches(6.2), h=Inches(2.7))
bl(sl, [
    "C0 (pure SVD) is fastest: 8 seconds. C1/C2/C4 take ~10 min each (per-user cosine similarity)",
    "C3 switching is faster (4.5 min): 57% of users skip the expensive content computation",
    "Per-user boxplot: most users score 0 (extreme sparsity), but C0/C2 show higher upper-quartile hits",
    "Runtime is acceptable for offline batch recommendation; online serving would need caching",
], l=Inches(0.5), t=Inches(4.5), w=Inches(12), h=Inches(2.5), sz=16)

# ======================== 23. CROSS-VERSION COMPARISON ========================
sl = ns(); tb(sl, "Cross-Version Comparison: A vs B vs C",
              "Which approach wins? It depends on what you optimize for")
cross = [
    ["Version", "Best Model", "NDCG@10", "Recall@10", "Coverage", "Best For"],
    ["A (Content)", "A4 Reranker", "0.00134", "0.00268", "45.30%", "Cold-start, explainability"],
    ["B (CF)", "B4 SVD", "0.01309", "0.02524", "3.26%", "Maximum precision"],
    ["C (Hybrid)", "C2 Weighted", "0.01194", "0.02200", "17.31%", "Balanced trade-off"],
]
tbl(sl, cross, Inches(0.3), Inches(1.5), [Inches(1.8), Inches(2), Inches(1.5), Inches(1.5), Inches(1.5), Inches(2.8)])
im(sl, os.path.join(C_FIG, "cross_version_metrics_comparison.png"), Inches(0.3), Inches(3.8), w=Inches(6.3), h=Inches(3.5))
im(sl, os.path.join(C_FIG, "cross_version_ndcg_vs_coverage.png"), Inches(6.8), Inches(3.8), w=Inches(6.2), h=Inches(3.5))

# ======================== 24. ALL MODELS + COVERAGE ========================
sl = ns(); tb(sl, "All Models Ranking & Coverage Trade-off")
im(sl, os.path.join(C_FIG, "all_models_ndcg_comparison.png"), Inches(0.3), Inches(1.5), w=Inches(6.3), h=Inches(5.5))
im(sl, os.path.join(C_FIG, "cross_version_coverage_comparison.png"), Inches(6.8), Inches(1.5), w=Inches(6.2), h=Inches(3))
bl(sl, [
    "The precision-coverage trade-off is the",
    "central design decision:",
    "",
    "  A: high coverage, low precision",
    "  B: high precision, low coverage",
    "  C: controllable balance between both",
], l=Inches(7), t=Inches(4.7), w=Inches(5.5), h=Inches(2.5), sz=16)

# ======================== 25. KEY FINDINGS ========================
sl = ns(); tb(sl, "Key Findings")
bl(sl, [
    "Finding 1: CF significantly outperforms content-based methods (10x NDCG improvement)",
    "    Interaction patterns are far more predictive than metadata similarity alone",
    "",
    "Finding 2: Pure CF suffers from critically low coverage (3.26%)",
    "    It creates an echo chamber, recommending only well-known items",
    "",
    "Finding 3: Hybrid methods successfully balance precision and coverage",
    "    C2: only 9% NDCG loss for 5x coverage gain. A strong real-world trade-off",
    "",
    "Finding 4: Content-based methods are essential for cold-start scenarios",
    "    C3 switching hybrid uses content for 43% of users with sparse history",
    "",
    "Finding 5: All metrics are low in absolute terms due to extreme sparsity (99.92%)",
    "    This is typical for large-scale recommendation; relative improvements matter more",
], sz=18)

# ======================== 26. PRACTICAL IMPLICATIONS ========================
sl = ns(); tb(sl, "Practical Implications")
bl(sl, [
    "For a real recipe recommendation application:",
    "",
    "    Homepage 'Top Picks for You' (precision-focused):",
    "        Use B4/C0 (SVD) for users with >= 10 interactions",
    "",
    "    'Explore New Recipes' (diversity-focused):",
    "        Use C3/C4 hybrid for broader recipe diversity",
    "",
    "    New user onboarding (cold-start):",
    "        Use C1/A4 content-based until enough ratings are collected",
    "        Prompt users to rate 5-10 recipes to bootstrap CF",
    "",
    "    Dietary constraints (practical filters):",
    "        Apply post-filters for prep time, ingredients, nutrition",
    "        Content methods naturally support constraint-based filtering",
    "",
    "    Recommended default: C2 Weighted Hybrid (alpha=0.7)",
], sz=18)

# ======================== 27. CONCLUSION ========================
sl = ns(); tb(sl, "Conclusion & Future Work")
bl(sl, [
    "Conclusion:",
    "    Best overall system: C2 Weighted Hybrid (alpha=0.7)",
    "    NDCG = 0.01194 (91% of pure CF), Coverage = 17.3% (5x pure CF)",
    "    Lightweight, interpretable, and practical for real deployment",
    "",
    "Limitations:",
    "    Extreme sparsity limits absolute performance across all methods",
    "    Single temporal split; k-fold cross-validation would improve robustness",
    "    No deep learning models compared",
    "",
    "Future Work:",
    "    Neural collaborative filtering (NCF) for better latent representations",
    "    LLM-based explanation layer ('We recommend this because...')",
    "    Online A/B testing to validate offline improvements with real users",
    "    Incorporate recipe images and cooking step sequences as features",
], sz=17)

# ======================== 28. THANK YOU ========================
sl = ns()
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg.fill.solid(); bg.fill.fore_color.rgb = DARK; bg.line.fill.background()
tx(sl, Inches(1), Inches(2.2), Inches(11), Inches(1.5), "Thank You", 54, True, WHITE, PP_ALIGN.CENTER)
tx(sl, Inches(1), Inches(3.8), Inches(11), Inches(1), "Questions?", 32, False, SUB, PP_ALIGN.CENTER)
tx(sl, Inches(1), Inches(5.5), Inches(11), Inches(0.6),
   "CS 513  |  NutriMatch  |  Hybrid Recipe Recommendation System", 18, False, MUT, PP_ALIGN.CENTER)

# ======================== SAVE ========================
prs.save(OUTPUT_PATH)
print(f"Saved: {os.path.abspath(OUTPUT_PATH)}")
print(f"Total slides: {len(prs.slides)}")
